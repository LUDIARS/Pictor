#!/usr/bin/env python3
"""Generate demo_guide.docx and demo_guide.pptx from content."""

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))


def create_docx():
    doc = Document()

    # Title
    title = doc.add_heading("Pictor Demo Guide", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(
        "Pictor が提供する 3 つのデモアプリケーションの解説資料です。"
        "各デモはそれぞれ独立したビルドターゲットとして構成されており、"
        "レンダリングパイプラインの異なる側面を実演します。"
    )

    # ---- Demo 1 ----
    doc.add_heading("1. Vulkan Window Demo (pictor_demo)", level=1)
    doc.add_heading("概要", level=2)
    doc.add_paragraph(
        "最もシンプルなデモ。GLFW ウィンドウを開き、Vulkan スワップチェーンを初期化して "
        "クリアカラーアニメーション（色相サイクル）を描画します。"
    )

    doc.add_heading("主な機能", level=2)
    table1 = doc.add_table(rows=5, cols=2, style="Light Grid Accent 1")
    headers1 = ["機能", "詳細"]
    data1 = [
        ("ウィンドウ", "GLFW 1280x720, VSync ON"),
        ("レンダリング", "クリアカラー HSV アニメーション (10秒で1周)"),
        ("Pictor 連携", "1000 個のテストオブジェクト登録"),
        ("バリデーション", "Vulkan Validation Layer 有効"),
    ]
    for j, h in enumerate(headers1):
        table1.rows[0].cells[j].text = h
    for i, (k, v) in enumerate(data1):
        table1.rows[i + 1].cells[0].text = k
        table1.rows[i + 1].cells[1].text = v

    doc.add_heading("ソースファイル", level=2)
    doc.add_paragraph("demo/main.cpp (約210行)", style="List Bullet")

    # ---- Demo 2 ----
    doc.add_heading("2. 1M Spheres Benchmark (pictor_benchmark)", level=1)
    doc.add_heading("概要", level=2)
    doc.add_paragraph(
        "100万個の球体を3D リサージュ曲線上で動かすストレステスト。"
        "GPU Compute Shader による更新と CPU 並列更新の両モードをサポートし、"
        "データパイプラインの性能を計測します。"
    )

    doc.add_heading("主な機能", level=2)
    table2 = doc.add_table(rows=6, cols=2, style="Light Grid Accent 1")
    headers2 = ["機能", "詳細"]
    data2 = [
        ("オブジェクト数", "1,000,000 (100x100x100 グリッド)"),
        ("更新モード", "GPU Compute (デフォルト) / CPU (--cpu)"),
        ("レンダリング", "SimpleRenderer (icosphere インスタンス描画)"),
        ("プロファイラ", "JSON/CSV/Chrome Tracing 出力"),
        ("ヘッドレス", "--headless でウィンドウなし実行"),
    ]
    for j, h in enumerate(headers2):
        table2.rows[0].cells[j].text = h
    for i, (k, v) in enumerate(data2):
        table2.rows[i + 1].cells[0].text = k
        table2.rows[i + 1].cells[1].text = v

    doc.add_heading("ソースファイル", level=2)
    doc.add_paragraph("benchmark/main.cpp (約380行)", style="List Bullet")

    # ---- Demo 3 ----
    doc.add_heading("3. PBR Graphics Demo (pictor_graphics_demo)", level=1)
    doc.add_heading("概要", level=2)
    doc.add_paragraph(
        "PBR (Physically Based Rendering)、シャドウマッピング、"
        "グローバルイルミネーション (GI) の統合動作をデモンストレーションするシーン。"
        "Cook-Torrance BRDF によるメタリックマテリアル、PCSS ソフトシャドウ、"
        "SSAO、ライトマップベイクを実演します。"
    )

    doc.add_heading("シーン構成", level=2)
    table3 = doc.add_table(rows=7, cols=4, style="Light Grid Accent 1")
    scene_headers = ["オブジェクト", "マテリアル", "配置", "種別"]
    scene_data = [
        ("メタリックキューブ", "Chrome (metalness=0.95)", "中央 Y=1.5", "Static"),
        ("地面", "Concrete (roughness=0.8)", "Y=0, 30x30", "Static"),
        ("銅球", "Copper (metalness=0.9)", "右前方", "Static"),
        ("金球", "Gold (metalness=0.95)", "左前方", "Static"),
        ("チタン球", "Titanium (metalness=0.85)", "後方", "Static"),
        ("ダイナミック球", "Teal Emissive", "円運動", "Dynamic"),
    ]
    for j, h in enumerate(scene_headers):
        table3.rows[0].cells[j].text = h
    for i, row in enumerate(scene_data):
        for j, val in enumerate(row):
            table3.rows[i + 1].cells[j].text = val

    doc.add_heading("ライティング", level=2)
    table4 = doc.add_table(rows=4, cols=4, style="Light Grid Accent 1")
    light_headers = ["ライト", "タイプ", "強度", "特徴"]
    light_data = [
        ("太陽光", "Directional", "2.5", "PCSS シャドウ 3カスケード"),
        ("スポットライト", "Spot", "4.0", "コーン20/35度, 周回運動"),
        ("アンビエント", "Hemisphere", "1.0", "天空+地面バウンス"),
    ]
    for j, h in enumerate(light_headers):
        table4.rows[0].cells[j].text = h
    for i, row in enumerate(light_data):
        for j, val in enumerate(row):
            table4.rows[i + 1].cells[j].text = val

    doc.add_heading("GI システム", level=2)
    doc.add_paragraph("シャドウ: CSM 3カスケード, 2048解像度, PCSS フィルタ", style="List Bullet")
    doc.add_paragraph("SSAO: 32サンプル, radius=0.5, bilateral blur", style="List Bullet")
    doc.add_paragraph("ライトマップベイク: 静的オブジェクトに事前計算 (bake_static_gi → apply_bake)", style="List Bullet")

    doc.add_heading("PBR シェーダ (Cook-Torrance BRDF)", level=2)
    doc.add_paragraph("D (NDF): GGX 正規分布関数", style="List Bullet")
    doc.add_paragraph("G (Geometry): Smith's Schlick-GGX", style="List Bullet")
    doc.add_paragraph("F (Fresnel): Schlick 近似 (F0 = mix(0.04, albedo, metallic))", style="List Bullet")
    doc.add_paragraph("トーンマッピング: ACES Filmic + sRGB ガンマ補正", style="List Bullet")

    doc.add_heading("ソースファイル", level=2)
    doc.add_paragraph("graphics_demo/main.cpp — デモアプリケーション + PBRDemoRenderer", style="List Bullet")
    doc.add_paragraph("shaders/pbr_demo.vert — PBR 頂点シェーダ", style="List Bullet")
    doc.add_paragraph("shaders/pbr_demo.frag — PBR フラグメントシェーダ", style="List Bullet")

    # ---- Build Options ----
    doc.add_heading("ビルドオプション一覧", level=1)
    table5 = doc.add_table(rows=4, cols=3, style="Light Grid Accent 1")
    build_headers = ["CMake オプション", "デフォルト", "説明"]
    build_data = [
        ("PICTOR_BUILD_DEMO", "ON", "Vulkan Window Demo"),
        ("PICTOR_BUILD_BENCHMARK", "ON", "1M Spheres Benchmark"),
        ("PICTOR_BUILD_GRAPHICS_DEMO", "ON", "PBR Graphics Demo"),
    ]
    for j, h in enumerate(build_headers):
        table5.rows[0].cells[j].text = h
    for i, row in enumerate(build_data):
        for j, val in enumerate(row):
            table5.rows[i + 1].cells[j].text = val

    path = os.path.join(OUTPUT_DIR, "demo_guide.docx")
    doc.save(path)
    print(f"Created: {path}")


def create_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_content_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    # Slide 1: Title
    add_title_slide("Pictor Demo Guide", "PBR + Shadow + GI デモ解説資料")

    # Slide 2: Overview
    add_content_slide("デモ一覧", [
        "1. Vulkan Window Demo (pictor_demo)",
        "   - 最小構成のVulkan パイプライン動作確認",
        "2. 1M Spheres Benchmark (pictor_benchmark)",
        "   - 100万オブジェクトのストレステスト",
        "3. PBR Graphics Demo (pictor_graphics_demo)",
        "   - PBR + Shadow + GI 統合デモ",
    ])

    # Slide 3: Vulkan Window Demo
    add_content_slide("1. Vulkan Window Demo", [
        "GLFW ウィンドウ + Vulkan スワップチェーン初期化",
        "クリアカラー HSV アニメーション (10秒で1周)",
        "1000 個テストオブジェクト登録 + プロファイラ統計",
        "Vulkan Validation Layer 有効",
        "ソース: demo/main.cpp (約210行)",
    ])

    # Slide 4: 1M Benchmark
    add_content_slide("2. 1M Spheres Benchmark", [
        "1,000,000 球体を 3D リサージュ曲線で移動",
        "GPU Compute (デフォルト) / CPU (--cpu) モード",
        "SimpleRenderer (icosphere instanced draw)",
        "プロファイラ出力: JSON / CSV / Chrome Tracing",
        "CPU フレーム時間目標: < 0.5ms",
        "ソース: benchmark/main.cpp (約380行)",
    ])

    # Slide 5: PBR Demo Overview
    add_content_slide("3. PBR Graphics Demo — 概要", [
        "Cook-Torrance BRDF (GGX + Smith + Schlick Fresnel)",
        "メタリックキューブ (metalness=0.95, roughness=0.15)",
        "地面 (concrete) + 3 静的メタリック球 (銅/金/チタン)",
        "1 動的球 (teal emissive, 円運動)",
        "ディレクショナルライト + スポットライト",
        "PCSS ソフトシャドウ + SSAO + ライトマップベイク",
    ])

    # Slide 6: Scene Layout
    add_content_slide("3. PBR Demo — シーン構成", [
        "メタリックキューブ: Chrome (中央, 回転アニメ, Static)",
        "地面: Concrete 30x30 (Y=0, Shadow Receiver)",
        "銅球: Copper metalness=0.9 (右前方, Static)",
        "金球: Gold metalness=0.95 (左前方, Static)",
        "チタン球: Titanium metalness=0.85 (後方, Static)",
        "ダイナミック球: Teal Emissive (円軌道, Dynamic)",
    ])

    # Slide 7: Lighting
    add_content_slide("3. PBR Demo — ライティング", [
        "太陽光 (Directional): 暖白, 強度 2.5",
        "  - PCSS シャドウ: 3カスケード CSM, 2048解像度",
        "  - コンタクトハードニングソフトシャドウ",
        "スポットライト (Spot): 暖色, 強度 4.0",
        "  - コーン: inner=20度, outer=35度",
        "  - シーン中心を照らしながら周回運動",
        "アンビエント: 半球ライティング (天空+地面バウンス)",
    ])

    # Slide 8: GI System
    add_content_slide("3. PBR Demo — GI システム", [
        "シャドウマッピング:",
        "  - CSM 3カスケード, 2048解像度, PCSS フィルタ",
        "  - depth bias + normal bias + slope-scale bias",
        "SSAO (Screen-Space Ambient Occlusion):",
        "  - 32サンプル, radius=0.5, bilateral blur",
        "ライトマップベイク:",
        "  - bake_static_gi() で静的オブジェクトを事前計算",
        "  - Shadow + AO + Lightmap を GPU にアップロード",
    ])

    # Slide 9: PBR Shader
    add_content_slide("3. PBR Demo — シェーダ詳細", [
        "Cook-Torrance Specular BRDF:",
        "  D: GGX Normal Distribution Function",
        "  G: Smith's Schlick-GGX Geometry",
        "  F: Schlick Fresnel (F0 = mix(0.04, albedo, metallic))",
        "Diffuse: Lambert (kD * albedo / PI)",
        "エネルギー保存: kD = (1-F) * (1-metallic)",
        "トーンマッピング: ACES Filmic + sRGB ガンマ補正",
    ])

    # Slide 10: Build
    add_content_slide("ビルド方法", [
        "cmake -DPICTOR_BUILD_DEMO=ON \\",
        "      -DPICTOR_BUILD_BENCHMARK=ON \\",
        "      -DPICTOR_BUILD_GRAPHICS_DEMO=ON ..",
        "make pictor_demo pictor_benchmark pictor_graphics_demo",
        "",
        "要件: CMake 3.20+, C++20, Vulkan SDK, GLFW3",
    ])

    path = os.path.join(OUTPUT_DIR, "demo_guide.pptx")
    prs.save(path)
    print(f"Created: {path}")


if __name__ == "__main__":
    create_docx()
    create_pptx()
